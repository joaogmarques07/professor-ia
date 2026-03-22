"""
Skill: Receptor de Conteúdo
Objetivo : Receber arquivos em diferentes formatos e extrair o texto bruto,
           pronto para ser processado pelas skills seguintes.
Suporta  : PDF, Word (.docx), PowerPoint (.pptx), Áudio/Vídeo (.mp3 .wav .m4a .mp4),
           Texto (.txt .md .csv), Excel (.xlsx), Web (.html .htm), Dados (.json .xml)
Restrições:
  - Documentos / Texto : máx 50 MB
  - Áudio / Vídeo      : máx 25 MB (limite da API Whisper)
  - Áudio requer OPENAI_API_KEY no .env
  - Esta skill NÃO processa nem interpreta o conteúdo — só extrai texto.
"""

import os
from dataclasses import dataclass, field
from io import BytesIO

# ─── Limites ────────────────────────────────────────────────────────────────
LIMITE_DOCUMENTO_MB = 50
LIMITE_AUDIO_MB = 25

FORMATOS_DOCUMENTO  = {".pdf", ".docx", ".pptx"}
FORMATOS_TEXTO      = {".txt", ".md", ".csv", ".json", ".xml"}
FORMATOS_WEB        = {".html", ".htm"}
FORMATOS_PLANILHA   = {".xlsx"}
FORMATOS_AUDIO      = {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".ogg", ".webm"}
FORMATOS_SUPORTADOS = (
    FORMATOS_DOCUMENTO | FORMATOS_TEXTO | FORMATOS_WEB | FORMATOS_PLANILHA | FORMATOS_AUDIO
)


# ─── Resultado ──────────────────────────────────────────────────────────────
@dataclass
class ConteudoExtraido:
    texto: str
    formato: str                        # "pdf" | "docx" | "pptx" | "audio"
    nome_arquivo: str
    metadados: dict = field(default_factory=dict)
    # metadados possíveis:
    #   paginas    (int)  — pdf, docx
    #   slides     (int)  — pptx
    #   duracao_seg (float) — audio


# ─── Ponto de entrada ────────────────────────────────────────────────────────
def extrair(conteudo_bytes: bytes, nome_arquivo: str) -> ConteudoExtraido:
    """
    Recebe os bytes do arquivo e o nome (usado para detectar o formato).
    Devolve um ConteudoExtraido com texto limpo e metadados.
    """
    nome = nome_arquivo.lower()
    extensao = "." + nome.rsplit(".", 1)[-1] if "." in nome else ""

    if extensao not in FORMATOS_SUPORTADOS:
        raise ValueError(
            f"Formato '{extensao}' não suportado. "
            f"Aceitos: {', '.join(sorted(FORMATOS_SUPORTADOS))}"
        )

    tamanho_mb = len(conteudo_bytes) / (1024 * 1024)
    limite = LIMITE_AUDIO_MB if extensao in FORMATOS_AUDIO else LIMITE_DOCUMENTO_MB

    if tamanho_mb > limite:
        raise ValueError(
            f"Arquivo muito grande: {tamanho_mb:.1f} MB. "
            f"Limite para {extensao}: {limite} MB."
        )

    if extensao == ".pdf":
        return _extrair_pdf(conteudo_bytes, nome_arquivo)
    elif extensao == ".docx":
        return _extrair_docx(conteudo_bytes, nome_arquivo)
    elif extensao == ".pptx":
        return _extrair_pptx(conteudo_bytes, nome_arquivo)
    elif extensao == ".xlsx":
        return _extrair_xlsx(conteudo_bytes, nome_arquivo)
    elif extensao in FORMATOS_WEB:
        return _extrair_html(conteudo_bytes, nome_arquivo)
    elif extensao in FORMATOS_TEXTO:
        return _extrair_texto_puro(conteudo_bytes, nome_arquivo, extensao)
    else:
        return _transcrever_audio(conteudo_bytes, nome_arquivo, extensao)


# ─── Extratores ──────────────────────────────────────────────────────────────
def _extrair_pdf(conteudo_bytes: bytes, nome: str) -> ConteudoExtraido:
    import pdfplumber

    paginas_texto = []
    with pdfplumber.open(BytesIO(conteudo_bytes)) as pdf:
        total = len(pdf.pages)
        for pagina in pdf.pages:
            t = pagina.extract_text()
            if t:
                paginas_texto.append(t.strip())

    if not paginas_texto:
        raise ValueError("Não foi possível extrair texto do PDF. O arquivo pode ser uma imagem escaneada.")

    return ConteudoExtraido(
        texto="\n\n".join(paginas_texto),
        formato="pdf",
        nome_arquivo=nome,
        metadados={"paginas": total},
    )


def _extrair_docx(conteudo_bytes: bytes, nome: str) -> ConteudoExtraido:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(BytesIO(conteudo_bytes))
    blocos = []

    for paragrafo in doc.paragraphs:
        t = paragrafo.text.strip()
        if t:
            blocos.append(t)

    # Extrair tabelas também
    for tabela in doc.tables:
        for linha in tabela.rows:
            celulas = [c.text.strip() for c in linha.cells if c.text.strip()]
            if celulas:
                blocos.append(" | ".join(celulas))

    if not blocos:
        raise ValueError("Não foi possível extrair texto do arquivo Word.")

    return ConteudoExtraido(
        texto="\n\n".join(blocos),
        formato="docx",
        nome_arquivo=nome,
        metadados={"paginas": len(doc.paragraphs)},
    )


def _extrair_pptx(conteudo_bytes: bytes, nome: str) -> ConteudoExtraido:
    from pptx import Presentation

    prs = Presentation(BytesIO(conteudo_bytes))
    slides_texto = []

    for i, slide in enumerate(prs.slides, start=1):
        partes = []

        for forma in slide.shapes:
            if forma.has_text_frame:
                texto = forma.text_frame.text.strip()
                if texto:
                    partes.append(texto)

        # Notas do apresentador
        if slide.has_notes_slide:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                partes.append(f"[Notas do slide {i}]: {notas}")

        if partes:
            slides_texto.append(f"--- Slide {i} ---\n" + "\n".join(partes))

    if not slides_texto:
        raise ValueError("Não foi possível extrair texto da apresentação.")

    return ConteudoExtraido(
        texto="\n\n".join(slides_texto),
        formato="pptx",
        nome_arquivo=nome,
        metadados={"slides": len(prs.slides)},
    )


def _extrair_texto_puro(conteudo_bytes: bytes, nome: str, extensao: str) -> ConteudoExtraido:
    """TXT, MD, CSV, JSON, XML — lê diretamente como texto."""
    texto = conteudo_bytes.decode("utf-8", errors="replace").strip()
    if not texto:
        raise ValueError("O arquivo está vazio.")
    return ConteudoExtraido(
        texto=texto,
        formato=extensao.lstrip("."),
        nome_arquivo=nome,
        metadados={"linhas": texto.count("\n") + 1},
    )


def _extrair_html(conteudo_bytes: bytes, nome: str) -> ConteudoExtraido:
    """HTML — remove tags e extrai só o texto visível."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(conteudo_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "head"]):
        tag.decompose()

    texto = soup.get_text(separator="\n").strip()
    linhas = [l.strip() for l in texto.splitlines() if l.strip()]
    texto_limpo = "\n".join(linhas)

    if not texto_limpo:
        raise ValueError("Não foi possível extrair texto do HTML.")

    return ConteudoExtraido(
        texto=texto_limpo,
        formato="html",
        nome_arquivo=nome,
        metadados={"linhas": len(linhas)},
    )


def _extrair_xlsx(conteudo_bytes: bytes, nome: str) -> ConteudoExtraido:
    """Excel (.xlsx) — extrai todas as planilhas como texto tabular."""
    import openpyxl

    wb = openpyxl.load_workbook(BytesIO(conteudo_bytes), data_only=True)
    blocos = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        linhas = []
        for row in ws.iter_rows(values_only=True):
            celulas = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in celulas):
                linhas.append(" | ".join(celulas))
        if linhas:
            blocos.append(f"=== Planilha: {sheet_name} ===\n" + "\n".join(linhas))

    if not blocos:
        raise ValueError("Não foi possível extrair dados da planilha.")

    return ConteudoExtraido(
        texto="\n\n".join(blocos),
        formato="xlsx",
        nome_arquivo=nome,
        metadados={"planilhas": len(wb.sheetnames)},
    )


def _transcrever_audio(conteudo_bytes: bytes, nome: str, extensao: str) -> ConteudoExtraido:
    """Transcreve áudio/vídeo usando a API Whisper da OpenAI."""
    from openai import OpenAI

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise EnvironmentError(
            "OPENAI_API_KEY não encontrada no .env. "
            "Necessária para transcrição de áudio."
        )

    client = OpenAI(api_key=api_key)

    # Whisper exige um nome de arquivo com extensão válida
    arquivo_virtual = (nome, BytesIO(conteudo_bytes), f"audio/{extensao.lstrip('.')}")

    transcricao = client.audio.transcriptions.create(
        model="whisper-1",
        file=arquivo_virtual,
        response_format="verbose_json",
        language="pt",
    )

    texto = transcricao.text.strip()
    duracao = getattr(transcricao, "duration", None)

    if not texto:
        raise ValueError("Não foi possível transcrever o áudio. O arquivo pode estar vazio ou silencioso.")

    return ConteudoExtraido(
        texto=texto,
        formato="audio",
        nome_arquivo=nome,
        metadados={"duracao_seg": round(duracao, 1) if duracao else None},
    )
